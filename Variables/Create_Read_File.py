from pptx import Presentation
import win32com.client

def create_new_file(file_name, content):
	file = open(file_name, "w+")
	file.write(content)
	file = open(file_name, "r")
	file_Content = file.read()
	return file_Content
	
def convert_file_into_pptx(file_name):	
	PptApp = win32com.client.Dispatch("Powerpoint.Application")
	PptApp.Visible = True
	PPtPresentation = PptApp.Presentations.Open(file_name)
	PPtPresentation.SaveAs(file_name+"x", 24)
	PPtPresentation.close()
	PptApp.Quit()
	
def open_and_read_file(file_name):
	file = open(file_name,"rb")
	prs = Presentation(file)
	fileDict = {}
	for slide in prs.slides:
		slideList = []
		for shape in slide.shapes:			
			if shape.has_text_frame:
				if shape.text=='Chart':
					continue
				if shape.text=='Table':
					continue
				if shape.text=='Photo':
					continue
				slideList.append(shape.text)
		print(slideList)		
		if slideList==[]:
			continue
		else:
			fileDict[slideList[0]] = slideList[1]
		print(fileDict)	
	return fileDict			